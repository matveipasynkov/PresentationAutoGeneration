import os
import requests
import subprocess
import speech_recognition as sr
from pptx import Presentation
from pptx.dml.color import RGBColor
import re
import time

# Настройки
PPTX_FILE = "auto_presentation.pptx"
OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL_NAME = "llama3.2" # Правильное имя модели

def recognize_speech():
    """Распознаёт речь через микрофон, возвращает текст или пустую строку."""
    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("\n🎤 Говорите... (скажите 'стоп' для выхода)")
        try:
            r.adjust_for_ambient_noise(source, duration=0.5)
            audio = r.listen(source, phrase_time_limit=15)
        except Exception as e:
            print(f"Ошибка записи: {e}")
            return ""
            
    try:
        text = r.recognize_google(audio, language="ru-RU")
        print(f"\n🔊 Распознано: {text}")
        return text
    except sr.UnknownValueError:
        print("Google Speech Recognition не смогла распознать речь")
        return ""
    except Exception as e:
        print(f"Ошибка распознавания: {e}")
        return ""

def generate_slide_data(text):
    """Генерирует заголовок, контент и дизайн для слайда."""
    title_response = ""
    content_lines = []
    design_suggestions = "Стандартный дизайн"
    
    if not text:
        print("Получен пустой текст для генерации")
        return "", [], design_suggestions
    
    try:
        # --- 1. Генерация заголовка ---
        title_prompt = f"""<|begin_of_text|><|start_header_id|>user<|end_header_id|>
Выдели основную тему из текста (3-5 слов):
"{text}". Ответ дай только самой темой без пояснений.<|eot_id|><|start_header_id|>assistant<|end_header_id|>"""
        
        print("Запрашиваю заголовок...")
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": MODEL_NAME,
                "prompt": title_prompt,
                "stream": False,
                "options": {"temperature": 0.3}
            },
            timeout=30
        )
        
        # Проверяем статус ответа
        if response.status_code != 200:
            print(f"Ошибка API (код {response.status_code}): {response.text}")
            return "Ошибка API", ["Не удалось получить ответ от API"], design_suggestions
            
        title_data = response.json()
        title_response = title_data.get("response", "").strip().replace('"', '')
        
        if not title_response:
            print("Получен пустой заголовок от API")
            title_response = "Тема: " + text[:20]  # Создаем заголовок из начала текста
        
        print(f"Заголовок: {title_response}")
        
        # --- 2. Генерация контента ---
        content_prompt = f"""<|begin_of_text|><|start_header_id|>user<|end_header_id|>
Сгенерируй 3 ключевых пункта для слайда на тему "{title_response}"
на основе текста: "{text}". Формат: маркированный список, начни каждый пункт с "-" или "*".
Отвечай на русском языке.<|eot_id|><|start_header_id|>assistant<|end_header_id|>"""
        
        print("Запрашиваю контент...")
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": MODEL_NAME,
                "prompt": content_prompt,
                "stream": False,
                "options": {"temperature": 0.5}
            },
            timeout=30
        )
        
        # Проверяем статус ответа
        if response.status_code != 200:
            print(f"Ошибка API для контента (код {response.status_code})")
            content_lines = [f"Пункт о {title_response}" for _ in range(3)]
        else:
            content_text = response.json().get("response", "")
            
            # Улучшенный парсинг контента с поддержкой разных маркеров
            raw_lines = content_text.split("\n")
            content_lines = []
            
            for line in raw_lines:
                # Убираем маркеры списка и начальные пробелы
                clean_line = re.sub(r'^[\s*•\-–—]+\s*', '', line).strip()
                if clean_line:  # Добавляем только непустые строки
                    content_lines.append(clean_line)
            
            # Если не смогли распарсить контент, создаем шаблонные пункты
            if not content_lines:
                print("Не удалось распарсить контент из ответа API")
                content_lines = [f"Ключевой аспект {i+1} темы '{title_response}'" for i in range(3)]
        
        # --- 3. Генерация дизайна ---
        design_prompt = f"""<|begin_of_text|><|start_header_id|>user<|end_header_id|>
Создай элегантный, профессиональный дизайн для слайда PowerPoint на тему "{title_response}".

ВАЖНЫЕ ТРЕБОВАНИЯ К ДИЗАЙНУ:
1. ШРИФТЫ: Используй ТОЛЬКО мягкие, современные шрифты. Обязательно выбери один из: Montserrat (предпочтительно), Raleway, Open Sans Light, Roboto Light для заголовка. Для основного текста подойдет Open Sans или Lato.

2. ЦВЕТА: Используй ТОЛЬКО спокойные, пастельные тона. Предпочтительно:
   - Основной цвет: молочный, белый, светло-серый, очень светлый бежевый (#F5F5F5, #FFFFF0, #FAFAFA)
   - Акцентный цвет: приглушенный, неяркий оттенок (пыльно-розовый, светло-голубой, нежно-зеленый)
   - НИКАКИХ ярких и контрастных сочетаний, которые напрягают глаза!

3. ФОН: Минималистичный, без лишних элементов. Предпочтительно однотонный светлый.

Опиши дизайн строго в формате:
1. Цветовая палитра: основной #HEX (светлый/молочный), акцентный #HEX (приглушенный, мягкий)
2. Шрифты: заголовок Montserrat Light/Raleway (или подобный мягкий шрифт), текст Open Sans/Lato
3. Фон: краткое описание (минималистичный, однотонный)

Дай ответ только в этом формате на русском языке.<|eot_id|><|start_header_id|>assistant<|end_header_id|>"""
        
        print("Запрашиваю дизайн...")
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": MODEL_NAME,
                "prompt": design_prompt,
                "stream": False,
                "options": {"temperature": 0.6}
            },
            timeout=30
        )
        
        if response.status_code == 200:
            design_suggestions = response.json().get("response", "Стандартный дизайн").strip()
        else:
            print(f"Ошибка API для дизайна (код {response.status_code})")
        
        return title_response, content_lines, design_suggestions
        
    except requests.exceptions.Timeout:
        print("Превышено время ожидания ответа от API")
        return "Таймаут API", ["Сервер не ответил вовремя"], "Стандартный дизайн"
    except requests.exceptions.ConnectionError:
        print(f"Не удалось подключиться к {OLLAMA_URL}")
        return "Ошибка соединения", ["Проверьте работу Ollama"], "Стандартный дизайн"
    except Exception as e:
        print(f"Неожиданная ошибка при генерации: {e}")
        return "Ошибка", ["Технические проблемы при генерации"], "Стандартный дизайн"

def parse_hex_color(text):
    """Извлекает HEX-код цвета из текста."""
    if not text: return None
    match = re.search(r'#([A-Fa-f0-9]{6}|[A-Fa-f0-9]{3})\b', str(text))
    if match:
        hex_color = match.group(1)
        if len(hex_color) == 3:  # Если формат #RGB
            hex_color = "".join([c*2 for c in hex_color])  # Преобразуем в #RRGGBB
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        except ValueError:
            print(f"Некорректный HEX код: {hex_color}")
    return None

def parse_design_suggestions(suggestions_text):
    """Парсит текстовые предложения по дизайну."""
    design = {
        'main_color': None, 
        'accent_color': None,
        'title_font': None, 
        'text_font': None, 
        'background_idea': None
    }
    
    # Парсим цвета
    colors = re.findall(r'#([A-Fa-f0-9]{6}|[A-Fa-f0-9]{3})\b', suggestions_text)
    if len(colors) >= 1:
        main_hex = colors[0]
        if len(main_hex) == 3:
            main_hex = ''.join([c*2 for c in main_hex])
        try:
            design['main_color'] = RGBColor(
                int(main_hex[0:2], 16),
                int(main_hex[2:4], 16),
                int(main_hex[4:6], 16)
            )
        except (ValueError, IndexError):
            print(f"Ошибка при парсинге основного цвета: {main_hex}")
            
    if len(colors) >= 2:
        accent_hex = colors[1]
        if len(accent_hex) == 3:
            accent_hex = ''.join([c*2 for c in accent_hex])
        try:
            design['accent_color'] = RGBColor(
                int(accent_hex[0:2], 16),
                int(accent_hex[2:4], 16),
                int(accent_hex[4:6], 16)
            )
        except (ValueError, IndexError):
            print(f"Ошибка при парсинге акцентного цвета: {accent_hex}")
    
    # Парсим шрифты - исправленная часть
    for line in suggestions_text.split('\n'):
        line = line.strip()
        # Проверяем на наличие информации о шрифте заголовка
        if 'заголовок' in line.lower():
            title_match = re.search(r'заголовок\s*:\s*([^,;]+)', line, re.IGNORECASE)
            if title_match:
                design['title_font'] = title_match.group(1).strip()
        
        # Проверяем на наличие информации о шрифте текста
        if 'текст' in line.lower():
            text_match = re.search(r'текст\s*:\s*([^,;]+)', line, re.IGNORECASE)
            if text_match:
                design['text_font'] = text_match.group(1).strip()
    
    # Парсим идею фона
    bg_match = re.search(r'фон\s*:\s*(.+)', suggestions_text, re.IGNORECASE)
    if bg_match:
        design['background_idea'] = bg_match.group(1).strip()
    
    # Значения по умолчанию
    if design['main_color'] is None: design['main_color'] = RGBColor(240, 240, 240)  # Светло-серый
    if design['accent_color'] is None: design['accent_color'] = RGBColor(0, 0, 0)  # Черный
    
    # Значения по умолчанию
    if design['main_color'] is None: design['main_color'] = RGBColor(240, 240, 240)  # Светло-серый
    if design['accent_color'] is None: design['accent_color'] = RGBColor(0, 0, 0)  # Черный

    # НОВЫЙ КОД: Проверка контрастности акцентного цвета
    main_brightness = sum([design['main_color'][0], design['main_color'][1], design['main_color'][2]]) / 3
    if main_brightness > 180:  # Если фон очень светлый
        accent_brightness = sum([design['accent_color'][0], design['accent_color'][1], design['accent_color'][2]]) / 3
        if accent_brightness > 160:  # Если акцентный цвет тоже светлый
            # Заменяем на темно-серый (мягче чем чистый черный)
            design['accent_color'] = RGBColor(60, 60, 60)
            print("⚠️ Акцентный цвет изменен на темный для улучшения читаемости на светлом фоне")
    print(f"ℹ️ Распарсенный дизайн: {design}")
    return design

def create_slide(presentation, title, content, design_suggestions):
    """Создаёт слайд и применяет базовый дизайн."""
    try:
        # Создаем слайд
        slide_layout = presentation.slide_layouts[1]  # Макет "Заголовок и контент"
        slide = presentation.slides.add_slide(slide_layout)
    except IndexError:
        try:
            slide_layout = presentation.slide_layouts[0]  # Пробуем первый макет
            slide = presentation.slides.add_slide(slide_layout)
        except Exception as e:
            print(f"Не удалось создать слайд: {e}")
            return

    # Парсим дизайн
    design = parse_design_suggestions(design_suggestions)
    main_color = design['main_color']
    accent_color = design['accent_color']
    title_font = design['title_font']
    text_font = design['text_font']

    # 1. Устанавливаем фон
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = main_color
        # Исправлено: RGBColor это кортеж (r,g,b), а не объект с атрибутом .rgb
        print(f"🎨 Установлен цвет фона: #{main_color[0]:02x}{main_color[1]:02x}{main_color[2]:02x}")
    except Exception as e:
        print(f"Не удалось установить фон: {e}")

    # Определяем контрастный цвет текста на основе яркости фона
    # Исправлено: используем индексирование для доступа к компонентам цвета
    brightness = sum([main_color[0], main_color[1], main_color[2]]) / 3
    text_color = RGBColor(0, 0, 0) if brightness > 128 else RGBColor(255, 255, 255)
    print(f"ℹ️ Яркость фона: {brightness:.0f}, выбран цвет текста: #{text_color[0]:02x}{text_color[1]:02x}{text_color[2]:02x}")

    # 2. Устанавливаем заголовок
    try:
        title_shape = slide.shapes.title
        if title_shape and hasattr(title_shape, 'text_frame'):
            # Очищаем и устанавливаем текст
            tf = title_shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = title[:50]  # Ограничиваем длину
            
            # Применяем стиль
            run.font.color.rgb = accent_color
            if title_font: run.font.name = title_font
            
            # Исправлено: используем индексирование для доступа к компонентам цвета
            print(f"🎨 Стилизован заголовок: Цвет #{accent_color[0]:02x}{accent_color[1]:02x}{accent_color[2]:02x}, Шрифт {title_font or 'default'}")
        else:
            print("Не найден подходящий шейп для заголовка")
    except Exception as e:
        print(f"Ошибка при стилизации заголовка: {e}")

    # 3. Устанавливаем контент
    try:
        content_box = None
        # Ищем подходящий плейсхолдер контента
        for shape in slide.placeholders:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type != 1:  # Не заголовок
                content_box = shape
                break
        
        if content_box and hasattr(content_box, 'text_frame'):
            tf = content_box.text_frame
            tf.clear()
            
            for i, line in enumerate(content[:3]):  # Макс. 3 пункта
                p = tf.add_paragraph()
                run = p.add_run()
                # Просто добавляем текст без дополнительного маркера
                run.text = line[:100]  # Ограничиваем длину
                run.font.color.rgb = text_color
                if text_font: run.font.name = text_font
            
            print(f"🎨 Стилизован контент: Цвет #{text_color[0]:02x}{text_color[1]:02x}{text_color[2]:02x}, Шрифт {text_font or 'default'}")
        else:
            print("Не найден плейсхолдер для контента")
    except Exception as e:
        print(f"Ошибка при стилизации контента: {e}")


def refresh_powerpoint():
    """Закрывает, сохраняет, открывает презентацию и переходит к последнему слайду."""
    script = f'''
    set pptFile to POSIX file "{os.path.abspath(PPTX_FILE)}"
    set pptFileName to "{os.path.basename(PPTX_FILE)}"

    tell application "Microsoft PowerPoint"
        activate
        
        # Закрываем текущую если есть
        try
            if (count of presentations) > 0 then
                if name of active presentation is pptFileName then
                    close active presentation saving yes
                    delay 1
                end if
            end if
        end try
        
        # Открываем файл
        try
            open pptFile
            delay 2
            
            # Переход к последнему слайду
            tell application "System Events"
                tell process "Microsoft PowerPoint"
                    # Cmd+Option+Down = переход к последнему слайду в PowerPoint
                    key code 125 using {{command down, option down}}
                end tell
            end tell
        end try
    end tell
    '''
    
    try:
        subprocess.run(["osascript", "-e", script], check=True)
    except subprocess.CalledProcessError as e:
        print(f"Ошибка выполнения AppleScript: {e}")
    except Exception as e:
        print(f"Неожиданная ошибка при обновлении PowerPoint: {e}")

def main():
    # Создаём или загружаем существующую презентацию
    try:
        if os.path.exists(PPTX_FILE):
            prs = Presentation(PPTX_FILE)
            print(f"🟢 Загружена презентация: {PPTX_FILE}")
        else:
            prs = Presentation()
            print("🟢 Создана новая презентация")
            prs.save(PPTX_FILE)
    except Exception as e:
        print(f"Ошибка при открытии/создании презентации: {e}")
        return

    try:
        while True:
            text = recognize_speech().strip()
            if not text:
                continue
                
            if "стоп" in text.lower():
                print("Получена команда 'стоп'. Завершение работы...")
                break

            print("\n⏳ Генерация данных слайда...")
            title, content, design_suggestions = generate_slide_data(text)
            
            print(f"\n📄 Заголовок: {title}")
            print("📌 Контент:")
            for item in content:
                print(f"   • {item}")
            print("\n🎨 Предложения по дизайну (от Ollama):")
            print(design_suggestions)

            print("\n🛠️ Создание и стилизация слайда...")
            create_slide(prs, title, content, design_suggestions)

            print("\n💾 Сохранение файла...")
            try:
                prs.save(PPTX_FILE)
                print("   Файл сохранен.")
                print("🔄 Обновление PowerPoint...")
                refresh_powerpoint()
                print("✅ Слайд добавлен и обновлен!")
            except Exception as e:
                print(f"Ошибка при сохранении: {e}")
                time.sleep(2)
                try:
                    # Повторная попытка сохранения
                    prs.save(PPTX_FILE)
                    refresh_powerpoint()
                    print("✅ Повторное сохранение успешно!")
                except Exception as e2:
                    print(f"Ошибка при повторном сохранении: {e2}")
                    
    except KeyboardInterrupt:
        print("\nПрервано пользователем")
    finally:
        try:
            prs.save(PPTX_FILE)
            print(f"\n💾 Финальное сохранение презентации как: {PPTX_FILE}")
        except Exception as e:
            print(f"Ошибка при финальном сохранении: {e}")

if __name__ == "__main__":
    # Проверка соединения с Ollama
    try:
        response = requests.get(OLLAMA_URL.replace("/api/generate", "/"), timeout=5)
        print(f"Проверка Ollama: {OLLAMA_URL.replace('/api/generate', '/')}...")
        if response.status_code == 200:
            print("🟢 Соединение с Ollama ОК.")
            main()
        else:
            print(f"🟠 Ollama ответила со статусом {response.status_code}. Проверьте сервер.")
    except requests.exceptions.ConnectionError:
        print("‼️ Не удалось подключиться к Ollama. Убедитесь, что сервер запущен.")
    except Exception as e:
        print(f"Ошибка при проверке соединения с Ollama: {e}")
